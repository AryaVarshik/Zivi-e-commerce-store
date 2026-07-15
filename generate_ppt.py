"""Generate a project presentation PPT for Zivi."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Brand colors (orange & white theme)
ORANGE = RGBColor(0xE8, 0x5D, 0x04)
DARK_ORANGE = RGBColor(0xDC, 0x2F, 0x02)
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xF9, 0xFA)
SOFT_ORANGE = RGBColor(0xFF, 0xF0, 0xE6)
MUTED = RGBColor(0x5C, 0x5C, 0x6E)
ACCENT = RGBColor(0xFB, 0x85, 0x00)


def set_run(run, size=18, bold=False, color=NAVY, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_round_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=NAVY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        p.level = 0
        run = p.add_run()
        run.text = "•  " + item
        set_run(run, size=size, color=color)
    return box


def footer(slide, page, total=9):
    add_textbox(
        slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
        "Zivi  ·  E-Commerce Web Application", size=11, color=MUTED
    )
    add_textbox(
        slide, Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.3),
        f"{page}/{total}", size=11, color=MUTED, align=PP_ALIGN.RIGHT
    )


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 9

    # ---- 1. Title ----
    s = prs.slides.add_slide(blank)
    add_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.25), Inches(7.5), ORANGE)
    add_textbox(s, Inches(1), Inches(2.0), Inches(11), Inches(1.2), "ZIVI", size=72, bold=True, color=ORANGE)
    add_textbox(s, Inches(1), Inches(3.2), Inches(11), Inches(0.6), "Premium Full-Stack E-Commerce Platform", size=28, color=WHITE)
    add_textbox(
        s, Inches(1), Inches(4.0), Inches(10), Inches(0.8),
        "A modern orange-and-white shopping experience built with Django,\nresponsive UI, AJAX cart, wishlist, and multi-step checkout.",
        size=16, color=RGBColor(0xB8, 0xB8, 0xC8)
    )
    add_textbox(s, Inches(1), Inches(5.5), Inches(10), Inches(0.4), "Project Presentation", size=14, color=ACCENT)

    # ---- 2. Agenda ----
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), ORANGE)
    add_textbox(s, Inches(0.7), Inches(0.5), Inches(10), Inches(0.6), "Agenda", size=36, bold=True, color=NAVY)
    topics = [
        "What is Zivi?",
        "Key Features",
        "Tech Stack",
        "Architecture & Structure",
        "User Journey",
        "Setup & Demo",
        "Future Roadmap",
        "Summary",
    ]
    for i, t in enumerate(topics):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 6)
        y = Inches(1.5 + row * 1.1)
        card = add_round_rect(s, x, y, Inches(5.5), Inches(0.9), SOFT_ORANGE)
        add_textbox(s, x + Inches(0.25), y + Inches(0.22), Inches(0.6), Inches(0.5), f"{i+1:02d}", size=22, bold=True, color=ORANGE)
        add_textbox(s, x + Inches(1.0), y + Inches(0.28), Inches(4), Inches(0.5), t, size=18, bold=True, color=NAVY)
    footer(s, 2, total)

    # ---- 3. What is Zivi ----
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), ORANGE)
    add_textbox(s, Inches(0.7), Inches(0.45), Inches(11), Inches(0.6), "What is Zivi?", size=36, bold=True, color=NAVY)
    add_textbox(
        s, Inches(0.7), Inches(1.25), Inches(12), Inches(1.0),
        "Zivi is a premium, fully responsive e-commerce web application with a polished "
        "orange-and-white design, dynamic product discovery, and a complete shopping flow "
        "from browse → cart → checkout.",
        size=17, color=MUTED
    )
    highlights = [
        ("240+", "Products seeded\nacross 12 categories"),
        ("Full Stack", "Django backend +\nvanilla HTML/CSS/JS"),
        ("AJAX Cart", "Instant totals without\npage reloads"),
        ("Themes", "Persistent light &\ndark mode support"),
    ]
    for i, (title, sub) in enumerate(highlights):
        x = Inches(0.7 + i * 3.15)
        add_round_rect(s, x, Inches(2.6), Inches(2.95), Inches(3.5), SOFT_ORANGE)
        add_textbox(s, x + Inches(0.2), Inches(3.2), Inches(2.55), Inches(0.8), title, size=26, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.2), Inches(4.2), Inches(2.55), Inches(1.2), sub, size=14, color=NAVY, align=PP_ALIGN.CENTER)
    footer(s, 3, total)

    # ---- 4. Key Features ----
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), ORANGE)
    add_textbox(s, Inches(0.7), Inches(0.4), Inches(11), Inches(0.55), "Key Features", size=34, bold=True, color=NAVY)
    left_items = [
        "Interactive 10-slide banner carousel",
        "Advanced filters (price, brand, rating, stock)",
        "AJAX shopping cart with live totals",
        "Wishlist with move-to-cart actions",
        "Multi-step checkout with gift options",
    ]
    right_items = [
        "Payment choices: COD, UPI, Cards, Net Banking",
        "Django auth for cart & wishlist safety",
        "Light / dark mode (localStorage)",
        "Reviews carousel + product review forms",
        "Flash sales page with countdown",
    ]
    add_round_rect(s, Inches(0.6), Inches(1.2), Inches(5.9), Inches(5.4), LIGHT_GRAY)
    add_round_rect(s, Inches(6.8), Inches(1.2), Inches(5.9), Inches(5.4), SOFT_ORANGE)
    add_textbox(s, Inches(0.9), Inches(1.45), Inches(5), Inches(0.4), "Shopping Experience", size=18, bold=True, color=ORANGE)
    add_bullets(s, Inches(0.9), Inches(2.1), Inches(5.3), Inches(4.2), left_items, size=15)
    add_textbox(s, Inches(7.1), Inches(1.45), Inches(5), Inches(0.4), "Account & Engagement", size=18, bold=True, color=DARK_ORANGE)
    add_bullets(s, Inches(7.1), Inches(2.1), Inches(5.3), Inches(4.2), right_items, size=15)
    footer(s, 4, total)

    # ---- 5. Tech Stack ----
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), ORANGE)
    add_textbox(s, Inches(0.7), Inches(0.4), Inches(11), Inches(0.55), "Tech Stack", size=34, bold=True, color=NAVY)

    layers = [
        ("Frontend", ORANGE, ["HTML5", "CSS3", "Vanilla JavaScript", "Font Awesome"]),
        ("Backend", DARK_ORANGE, ["Django 6", "Python 3.13", "Django Templates", "AJAX APIs"]),
        ("Data & Deploy", ACCENT, ["SQLite (dev)", "PostgreSQL ready", "Gunicorn", "WhiteNoise"]),
    ]
    for i, (title, color, items) in enumerate(layers):
        x = Inches(0.7 + i * 4.15)
        add_round_rect(s, x, Inches(1.3), Inches(3.9), Inches(5.2), LIGHT_GRAY)
        add_rect(s, x, Inches(1.3), Inches(3.9), Inches(0.7), color)
        add_textbox(s, x, Inches(1.4), Inches(3.9), Inches(0.5), title, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_bullets(s, x + Inches(0.35), Inches(2.4), Inches(3.3), Inches(3.5), items, size=16)
    footer(s, 5, total)

    # ---- 6. Architecture ----
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), ORANGE)
    add_textbox(s, Inches(0.7), Inches(0.4), Inches(11), Inches(0.55), "Architecture & Structure", size=34, bold=True, color=NAVY)

    boxes = [
        (0.7, "zivi_project/", "Settings, URLs,\nWSGI / ASGI"),
        (4.0, "store/", "Models, Views,\nTemplates, Static"),
        (7.3, "Data Layer", "Products, Cart,\nOrders, Reviews"),
        (10.5, "Bootstrap", "run.py + seed\n240+ products"),
    ]
    for x, title, desc in boxes:
        add_round_rect(s, Inches(x), Inches(1.5), Inches(2.9), Inches(2.4), SOFT_ORANGE)
        add_textbox(s, Inches(x + 0.15), Inches(1.75), Inches(2.6), Inches(0.5), title, size=16, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(x + 0.15), Inches(2.4), Inches(2.6), Inches(1.2), desc, size=14, color=NAVY, align=PP_ALIGN.CENTER)

    # arrows between boxes
    for i in range(3):
        ax = Inches(3.55 + i * 3.25)
        add_textbox(s, ax, Inches(2.4), Inches(0.5), Inches(0.4), "→", size=24, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(0.7), Inches(4.3), Inches(12), Inches(0.4), "Core pages", size=18, bold=True, color=NAVY)
    pages = "Home  ·  Product Detail  ·  Flash Sale  ·  Cart  ·  Wishlist  ·  Checkout  ·  Profile  ·  Settings  ·  Auth"
    add_round_rect(s, Inches(0.7), Inches(4.9), Inches(11.9), Inches(1.3), LIGHT_GRAY)
    add_textbox(s, Inches(0.95), Inches(5.25), Inches(11.4), Inches(0.7), pages, size=15, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s, 6, total)

    # ---- 7. User Journey ----
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), ORANGE)
    add_textbox(s, Inches(0.7), Inches(0.4), Inches(11), Inches(0.55), "User Journey", size=34, bold=True, color=NAVY)

    steps = [
        ("01", "Browse", "Explore banners,\ncategories & filters"),
        ("02", "Engage", "Wishlist items &\nread/write reviews"),
        ("03", "Cart", "AJAX quantity\nupdates & live totals"),
        ("04", "Checkout", "Address, gifts &\npayment method"),
        ("05", "Confirm", "Order success &\nprofile history"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        x = Inches(0.55 + i * 2.55)
        add_round_rect(s, x, Inches(1.6), Inches(2.35), Inches(4.4), LIGHT_GRAY if i % 2 == 0 else SOFT_ORANGE)
        add_textbox(s, x + Inches(0.15), Inches(2.0), Inches(2.05), Inches(0.5), num, size=28, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.15), Inches(2.9), Inches(2.05), Inches(0.5), title, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.15), Inches(3.7), Inches(2.05), Inches(1.4), desc, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s, 7, total)

    # ---- 8. Setup ----
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), ORANGE)
    add_textbox(s, Inches(0.7), Inches(0.4), Inches(11), Inches(0.55), "Setup & Demo", size=34, bold=True, color=NAVY)

    add_round_rect(s, Inches(0.7), Inches(1.3), Inches(7.2), Inches(5.2), NAVY)
    add_textbox(s, Inches(1.0), Inches(1.55), Inches(6.5), Inches(0.4), "Quick start", size=18, bold=True, color=ORANGE)
    commands = [
        "python -m venv venv",
        "pip install -r requirements.txt",
        "python run.py",
        "",
        "# run.py will:",
        "#  • run migrations",
        "#  • seed 240+ products",
        "#  • start http://127.0.0.1:8000/",
    ]
    for i, line in enumerate(commands):
        add_textbox(s, Inches(1.0), Inches(2.15 + i * 0.45), Inches(6.5), Inches(0.4), line, size=15, color=WHITE)

    add_round_rect(s, Inches(8.2), Inches(1.3), Inches(4.4), Inches(5.2), SOFT_ORANGE)
    add_textbox(s, Inches(8.5), Inches(1.7), Inches(3.9), Inches(0.5), "Admin Access", size=18, bold=True, color=DARK_ORANGE)
    add_bullets(
        s, Inches(8.5), Inches(2.5), Inches(3.9), Inches(3.5),
        [
            "URL: /admin/",
            "Username: admin",
            "Password: admin123",
            "Manage catalog & orders",
        ],
        size=15
    )
    footer(s, 8, total)

    # ---- 9. Roadmap + Thank you ----
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), ORANGE)
    add_textbox(s, Inches(0.7), Inches(0.4), Inches(11), Inches(0.55), "Future Roadmap", size=34, bold=True, color=NAVY)

    roadmap = [
        ("Cloud Images", "Connect product images to S3 / cloud storage"),
        ("Real Payments", "Integrate Stripe, PayPal, or UPI gateways"),
        ("Order Tracking", "Live delivery status + SMS/email alerts"),
    ]
    for i, (title, desc) in enumerate(roadmap):
        y = Inches(1.3 + i * 1.35)
        add_round_rect(s, Inches(0.7), y, Inches(12), Inches(1.15), LIGHT_GRAY if i % 2 == 0 else SOFT_ORANGE)
        add_textbox(s, Inches(1.0), y + Inches(0.2), Inches(11), Inches(0.4), title, size=18, bold=True, color=ORANGE)
        add_textbox(s, Inches(1.0), y + Inches(0.6), Inches(11), Inches(0.4), desc, size=14, color=NAVY)

    # closing strip
    add_rect(s, Inches(0), Inches(5.6), Inches(13.333), Inches(1.9), NAVY)
    add_textbox(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.5), "Thank You", size=32, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_textbox(
        s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.4),
        "Zivi — shop smarter, build better.", size=16, color=WHITE, align=PP_ALIGN.CENTER
    )

    out = r"c:\Users\aryav\Documents\antigravity\zivi\Zivi_Project_Presentation.pptx"
    prs.save(out)
    print(out)


if __name__ == "__main__":
    build()
